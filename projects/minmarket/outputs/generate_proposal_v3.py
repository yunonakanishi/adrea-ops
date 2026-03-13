"""
くらしのマーケット × アドレア 紙媒体プロモーション提案書 v3
README.mdのスライド構成（10枚）に基づいて生成
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── カラーパレット ───
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
BLUE = RGBColor(0x2E, 0x86, 0xC1)
LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
ACCENT = RGBColor(0xE7, 0x4C, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF3, 0xF4)
GREEN = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def set_text(shape, text, font_size=14, color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_textbox(slide, left, top, width, height, text, font_size=14, color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_multi_text(slide, left, top, width, height, lines):
    """lines: list of (text, font_size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, font_size, color, bold, alignment) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = alignment
        p.space_after = Pt(4)
    return tf


def slide_header(slide, title, subtitle=None):
    """共通ヘッダー：上部にNAVYバー + タイトル"""
    add_rect(slide, 0, 0, W, Inches(1.1), NAVY)
    add_textbox(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.6),
                title, font_size=28, color=WHITE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(0.65), Inches(11), Inches(0.4),
                    subtitle, font_size=14, color=RGBColor(0xAA, 0xCC, 0xEE))
    # ページ下部にアクセントライン
    add_rect(slide, 0, H - Inches(0.06), W, Inches(0.06), BLUE)


# ══════════════════════════════════════════════════
# Slide 1: 表紙
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, NAVY)

# アクセントライン
add_rect(slide, Inches(0.8), Inches(2.0), Inches(1.5), Inches(0.06), BLUE)

add_multi_text(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(3.0), [
    ("紙媒体プロモーションのご提案", 40, WHITE, True, PP_ALIGN.LEFT),
    ("", 12, WHITE, False, PP_ALIGN.LEFT),
    ("地方エリアでの「選ばれるプラットフォーム」を\n紙媒体で確立する", 22, RGBColor(0xAA, 0xCC, 0xEE), False, PP_ALIGN.LEFT),
])

add_multi_text(slide, Inches(0.8), Inches(5.5), Inches(6), Inches(1.5), [
    ("くらしのマーケット 様", 18, RGBColor(0xDD, 0xDD, 0xDD), False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("2026年3月  |  株式会社アドレア", 14, RGBColor(0x99, 0xAA, 0xBB), False, PP_ALIGN.LEFT),
])

add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), BLUE)


# ══════════════════════════════════════════════════
# Slide 2: くらしのマーケットの成長と次の成長レバー
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "くらしのマーケットの成長と次の成長レバー",
             "圧倒的な実績の先にある、次のステージ")

# 3つの実績カード
cards = [
    ("10万店+", "累計出店登録数", "国内最大級の出張訪問\nサービスプラットフォーム"),
    ("400+", "サービスカテゴリ数", "生活に関わるあらゆる\nサービスをカバー"),
    ("ROAS 150%", "ペイド広告の成果", "目標超過達成\n高い広告効率を実現"),
]
for i, (num, label, desc) in enumerate(cards):
    x = Inches(0.8 + i * 4.1)
    y = Inches(1.6)
    card = add_rounded_rect(slide, x, y, Inches(3.6), Inches(2.5), LIGHT_BLUE)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.0), Inches(0.8),
                num, font_size=36, color=BLUE, bold=True)
    add_textbox(slide, x + Inches(0.3), y + Inches(1.0), Inches(3.0), Inches(0.4),
                label, font_size=14, color=NAVY, bold=True)
    add_textbox(slide, x + Inches(0.3), y + Inches(1.5), Inches(3.0), Inches(0.9),
                desc, font_size=12, color=MID_GRAY)

# 矢印セクション
add_rect(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.06), BLUE)

add_multi_text(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.2), [
    ("▶  次の成長レバー", 20, NAVY, True, PP_ALIGN.LEFT),
    ("", 6, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("① 地方エリアの立ち上げ（北陸・四国を最優先、次いで九州）", 16, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("② 売り手（職人）の集客拡大 — 個人事業主が多く、従来の営業手法ではリーチが困難", 16, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("③ 15周年（2026年7月）を起点としたプロモーション強化", 16, DARK_GRAY, False, PP_ALIGN.LEFT),
])


# ══════════════════════════════════════════════════
# Slide 3: 地方エリアが持つ2つの構造的な壁
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "地方エリアが持つ2つの構造的な壁",
             "デジタルだけでは届きにくい領域がある")

# 壁1
box1 = add_rounded_rect(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.5), RGBColor(0xFD, 0xED, 0xEC))
add_multi_text(slide, Inches(1.2), Inches(1.8), Inches(4.8), Inches(4.0), [
    ("壁 ①", 14, ACCENT, True, PP_ALIGN.LEFT),
    ("ユーザーと職人の\n「ニワトリとタマゴ」問題", 22, NAVY, True, PP_ALIGN.LEFT),
    ("", 8, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("地方では職人が少ない", 14, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("　→ ユーザーが検索しても選択肢が少ない", 14, MID_GRAY, False, PP_ALIGN.LEFT),
    ("　→ ユーザーが離脱する", 14, MID_GRAY, False, PP_ALIGN.LEFT),
    ("　→ 職人にとって出店メリットが薄い", 14, MID_GRAY, False, PP_ALIGN.LEFT),
    ("", 8, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("ユーザーと職人、両サイドを同時に立ち上げる\n施策が必要", 14, ACCENT, True, PP_ALIGN.LEFT),
])

# 壁2
box2 = add_rounded_rect(slide, Inches(6.9), Inches(1.6), Inches(5.5), Inches(4.5), RGBColor(0xEB, 0xF5, 0xFB))
add_multi_text(slide, Inches(7.3), Inches(1.8), Inches(4.8), Inches(4.0), [
    ("壁 ②", 14, BLUE, True, PP_ALIGN.LEFT),
    ("デジタル広告の\n「地方の壁」", 22, NAVY, True, PP_ALIGN.LEFT),
    ("", 8, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("地方はデジタル広告の配信ボリュームが限定的", 14, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("　→ リスティング・ディスプレイの在庫が少ない", 14, MID_GRAY, False, PP_ALIGN.LEFT),
    ("　→ CPAが高騰しやすい", 14, MID_GRAY, False, PP_ALIGN.LEFT),
    ("　→ そもそもリーチできない層が存在する", 14, MID_GRAY, False, PP_ALIGN.LEFT),
    ("", 8, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("デジタル以外のチャネルで\n「面」を取る必要がある", 14, BLUE, True, PP_ALIGN.LEFT),
])


# ══════════════════════════════════════════════════
# Slide 4: 提案：テレビCMで認知 × 紙媒体で行動
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "提案：認知施策 × 紙媒体で行動喚起",
             "3つの柱で地方エリアの両サイド同時獲得を実現")

pillars = [
    ("01", "両サイド同時獲得", ACCENT,
     "ユーザー向け（表面）と\n職人向け（裏面）を\n1枚のチラシに集約\n\nユーザーと売り手の\n両方に同時リーチ"),
    ("02", "エリアピンポイント", BLUE,
     "町丁目単位のデータで\n配布エリアを精密設計\n\n北陸・四国の\n優先エリアに集中投下\nムダ打ちを最小化"),
    ("03", "紙媒体フルラインナップ", GREEN,
     "ポスティング\n新聞折込\nDM\nサンプリング\n\nエリア特性に応じた\n最適な組み合わせ"),
]

for i, (num, title, color, desc) in enumerate(pillars):
    x = Inches(0.8 + i * 4.1)
    y = Inches(1.5)
    # 番号サークル
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.1), y + Inches(0.2), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    set_text(circle, num, font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf = circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_textbox(slide, x + Inches(1.0), y + Inches(0.25), Inches(2.8), Inches(0.6),
                title, font_size=20, color=NAVY, bold=True)

    card = add_rounded_rect(slide, x, y + Inches(1.1), Inches(3.6), Inches(4.2), LIGHT_GRAY)
    add_textbox(slide, x + Inches(0.3), y + Inches(1.3), Inches(3.0), Inches(3.8),
                desc, font_size=14, color=DARK_GRAY)


# ══════════════════════════════════════════════════
# Slide 5: なぜアドレアか：貴社が得られる4つの価値
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "なぜアドレアか：貴社が得られる4つの価値")

values = [
    ("📊", "データに基づく配布設計",
     "町丁目単位のエリアDB（人口統計・世帯構成・住居形態）を保持。\n貴社DBとの掛け合わせで、最も反応が見込めるエリアを特定します。"),
    ("🔧", "企画から配布まで一気通貫",
     "クリエイティブ制作・印刷・配布・効果計測まで一括対応。\n複数ベンダーとの調整コストをゼロにします。"),
    ("📰", "紙媒体の最適な組み合わせ",
     "ポスティング・新聞折込・DM・サンプリング等の豊富な実績。\nエリア特性に応じたメディアミックスをご提案します。"),
    ("🔄", "改善が回り続ける仕組み",
     "フードデリバリーmenuで構築したPDCAフレームを転用。\nエリア別CVR分析 → 訴求軸改善 → 配布最適化のサイクルを回します。"),
]

for i, (icon, title, desc) in enumerate(values):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.5 + row * 2.8)

    card = add_rounded_rect(slide, x, y, Inches(5.8), Inches(2.4), LIGHT_GRAY)

    add_textbox(slide, x + Inches(0.3), y + Inches(0.2), Inches(0.6), Inches(0.5),
                icon, font_size=28, color=BLUE, bold=True)
    add_textbox(slide, x + Inches(1.0), y + Inches(0.25), Inches(4.5), Inches(0.5),
                title, font_size=18, color=NAVY, bold=True)
    add_textbox(slide, x + Inches(1.0), y + Inches(0.85), Inches(4.5), Inches(1.4),
                desc, font_size=13, color=MID_GRAY)


# ══════════════════════════════════════════════════
# Slide 6: プランの全体像：テスト → 15周年本番
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "プランの全体像：テスト → 15周年本番",
             "小さく検証し、15周年で一気に拡大する2段構え")

# テスト期
add_rect(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.6), BLUE)
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.6),
            "テスト期（4〜5月）", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

test_items = [
    "対象エリア：北陸（石川 or 富山）1〜2市",
    "配布部数：3〜5万部（ポスティング中心）",
    "想定予算：200〜400万円",
    "チラシ仕様：両面（表：ユーザー訴求 / 裏：職人募集）",
    "計測：プロモーションコード or 専用LP",
    "目的：エリア特性の把握、反応率のベンチマーク取得",
]
for i, item in enumerate(test_items):
    add_textbox(slide, Inches(1.2), Inches(2.4 + i * 0.45), Inches(4.8), Inches(0.4),
                f"•  {item}", font_size=13, color=DARK_GRAY)

# 本番期
add_rect(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.6), ACCENT)
add_textbox(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.6),
            "15周年本番（7月）", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

main_items = [
    "対象エリア：北陸全域 + 四国（テスト結果で調整）",
    "配布部数：20〜50万部（ポスティング + 新聞折込）",
    "想定予算：1,000〜2,000万円",
    "メディアミックス：ポスティング + 新聞折込 + DM",
    "15周年クリエイティブで認知 × 行動喚起を両立",
    "NOBROCK TV施策との連動（認知→行動の導線設計）",
]
for i, item in enumerate(main_items):
    add_textbox(slide, Inches(7.4), Inches(2.4 + i * 0.45), Inches(4.8), Inches(0.4),
                f"•  {item}", font_size=13, color=DARK_GRAY)

# 矢印
arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.1), Inches(3.5), Inches(0.8), Inches(0.5))
arrow.fill.solid()
arrow.fill.fore_color.rgb = NAVY
arrow.line.fill.background()

# 下部補足
add_rounded_rect(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.3), RGBColor(0xFE, 0xF9, 0xE7))
add_multi_text(slide, Inches(1.2), Inches(5.6), Inches(11.0), Inches(1.1), [
    ("💡 予算の考え方", 14, NAVY, True, PP_ALIGN.LEFT),
    ("テスト期は200〜400万円の小規模投資でリスクを抑えつつ検証。本番は結果を見て1,000〜2,000万円規模に拡大。", 13, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("月間広告予算（約6,000万円）の中で、新規チャネル開拓として段階的に投下するプランです。", 13, MID_GRAY, False, PP_ALIGN.LEFT),
])


# ══════════════════════════════════════════════════
# Slide 7: 配布設計：データの掛け合わせ × 両面チラシ
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "配布設計：データの掛け合わせ × 両面チラシ",
             "貴社DBとアドレアDBで最適エリアを特定し、1枚で両サイドにリーチ")

# DB掛け合わせ
add_rounded_rect(slide, Inches(0.8), Inches(1.5), Inches(3.6), Inches(2.0), LIGHT_BLUE)
add_multi_text(slide, Inches(1.1), Inches(1.6), Inches(3.0), Inches(1.8), [
    ("貴社DB", 16, BLUE, True, PP_ALIGN.CENTER),
    ("", 4, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("• 市区町村別の出店者数", 13, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("• カテゴリ別の需給バランス", 13, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("• エリア別のユーザー数", 13, DARK_GRAY, False, PP_ALIGN.LEFT),
])

# × マーク
add_textbox(slide, Inches(4.5), Inches(2.0), Inches(0.8), Inches(0.8),
            "×", font_size=32, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, Inches(5.4), Inches(1.5), Inches(3.6), Inches(2.0), RGBColor(0xE8, 0xF8, 0xF5))
add_multi_text(slide, Inches(5.7), Inches(1.6), Inches(3.0), Inches(1.8), [
    ("アドレアDB", 16, GREEN, True, PP_ALIGN.CENTER),
    ("", 4, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("• 町丁目単位の人口統計", 13, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("• 世帯構成・住居形態", 13, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("• menu実績のCVR傾向値", 13, DARK_GRAY, False, PP_ALIGN.LEFT),
])

# → 最適エリア
add_textbox(slide, Inches(9.1), Inches(2.0), Inches(0.8), Inches(0.8),
            "→", font_size=32, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, Inches(9.9), Inches(1.5), Inches(2.6), Inches(2.0), RGBColor(0xFD, 0xED, 0xEC))
add_multi_text(slide, Inches(10.1), Inches(1.6), Inches(2.2), Inches(1.8), [
    ("最適配布エリア", 16, ACCENT, True, PP_ALIGN.CENTER),
    ("", 4, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("職人不足 × 需要見込\nが高いエリアを\n町丁目単位で特定", 13, DARK_GRAY, False, PP_ALIGN.CENTER),
])

# 両面チラシ
add_rect(slide, Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.06), BLUE)
add_textbox(slide, Inches(0.8), Inches(4.2), Inches(4.0), Inches(0.5),
            "▶  両面チラシの構造", 18, NAVY, True)

# 表面
add_rounded_rect(slide, Inches(0.8), Inches(4.9), Inches(5.5), Inches(2.0), RGBColor(0xEB, 0xF5, 0xFB))
add_multi_text(slide, Inches(1.2), Inches(5.0), Inches(4.8), Inches(1.8), [
    ("【表面】ユーザー向け", 16, BLUE, True, PP_ALIGN.LEFT),
    ("", 4, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("「あなたの街のプロに頼める」訴求", 14, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("15周年キャンペーン特典（初回割引等）", 14, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("QRコード → 専用LP（計測用）", 14, DARK_GRAY, False, PP_ALIGN.LEFT),
])

# 裏面
add_rounded_rect(slide, Inches(7.0), Inches(4.9), Inches(5.5), Inches(2.0), RGBColor(0xFD, 0xED, 0xEC))
add_multi_text(slide, Inches(7.4), Inches(5.0), Inches(4.8), Inches(1.8), [
    ("【裏面】職人・事業者向け", 16, ACCENT, True, PP_ALIGN.LEFT),
    ("", 4, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("「初期費用ゼロ・月額ゼロで集客」訴求", 14, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("出店登録のメリット（手数料は成約時のみ）", 14, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("QRコード → 出店登録ページ（計測用）", 14, DARK_GRAY, False, PP_ALIGN.LEFT),
])


# ══════════════════════════════════════════════════
# Slide 8: スケジュール
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "スケジュール",
             "3月スタート → 7月15周年本番の全体タイムライン")

phases = [
    ("3月", "準備", RGBColor(0x85, 0x92, 0x9E),
     ["DB共有・エリア選定", "見積作成", "CR方向性合意"]),
    ("4月", "制作", BLUE,
     ["チラシデザイン制作", "印刷・配布手配", "計測設計確定"]),
    ("4〜5月", "テスト配布", GREEN,
     ["北陸1〜2市で配布", "3〜5万部", "効果計測・分析"]),
    ("6月", "改善・拡大準備", ORANGE,
     ["テスト結果分析", "CR改善・エリア拡大設計", "本番用制作"]),
    ("7月", "15周年本番", ACCENT,
     ["北陸全域+四国展開", "20〜50万部", "NOBROCK TV連動"]),
]

bar_y = Inches(2.0)
bar_h = Inches(0.8)
total_w = Inches(11.7)
start_x = Inches(0.8)

for i, (month, phase, color, items) in enumerate(phases):
    x = start_x + Emu(int(total_w * i / 5))
    w = Emu(int(total_w / 5)) - Inches(0.1)

    # フェーズバー
    bar = add_rounded_rect(slide, x, bar_y, w, bar_h, color)
    set_text(bar, f"{month}\n{phase}", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    bar.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    if len(bar.text_frame.paragraphs) > 1:
        bar.text_frame.paragraphs[1].alignment = PP_ALIGN.CENTER

    # タスク
    for j, item in enumerate(items):
        add_textbox(slide, x + Inches(0.1), Inches(3.1 + j * 0.45), w - Inches(0.1), Inches(0.4),
                    f"• {item}", font_size=12, color=DARK_GRAY)

# 下部にマイルストーン
add_rect(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.04), MID_GRAY)
milestones = [
    (0.0, "キックオフ"),
    (0.4, "テスト配布開始"),
    (0.6, "中間レビュー"),
    (0.8, "本番配布開始"),
]
for pos, label in milestones:
    mx = start_x + Emu(int(total_w * pos))
    # ドット
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, mx - Inches(0.08), Inches(4.94), Inches(0.16), Inches(0.16))
    dot.fill.solid()
    dot.fill.fore_color.rgb = NAVY
    dot.line.fill.background()
    add_textbox(slide, mx - Inches(0.8), Inches(5.15), Inches(1.6), Inches(0.4),
                label, font_size=11, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

# 北陸メディア参考情報
add_rounded_rect(slide, Inches(0.8), Inches(5.7), Inches(11.7), Inches(1.2), LIGHT_GRAY)
add_multi_text(slide, Inches(1.2), Inches(5.8), Inches(11.0), Inches(1.0), [
    ("📰 北陸エリアのメディア参考情報", 13, NAVY, True, PP_ALIGN.LEFT),
    ("富山：北日本新聞（約23万部・シェア67%）/ 石川：北國新聞（折込エリアマップ公開）/ ポスティング：カラフルカンパニー（戸別宅配型フリーペーパー運営）", 11, MID_GRAY, False, PP_ALIGN.LEFT),
    ("→ 新聞折込は地方ほど相対的にリーチ力が高い（全国的な減少傾向の中で中部は世帯あたり部数が高水準）", 11, MID_GRAY, False, PP_ALIGN.LEFT),
])


# ══════════════════════════════════════════════════
# Slide 9: ネクストステップ
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "ネクストステップ",
             "4つのアクションで具体化を進めます")

steps = [
    ("STEP 1", "DB共有", "貴社の市区町村別データ（出店者数・\nユーザー数）を共有いただき、\nアドレアDBと掛け合わせて\n最適配布エリアを設計します", BLUE),
    ("STEP 2", "見積作成", "北陸・四国のポスティング単価・\n配布カバー率を確認し、\nテスト配布の具体見積を\nご提出します", GREEN),
    ("STEP 3", "CR方向性合意", "両面チラシの訴求軸・\nデザイン方向性について\n合意いただき、\n制作に着手します", ORANGE),
    ("STEP 4", "計測設計確定", "プロモーションコード or 専用LP\nによる計測方法を確定し、\n効果測定の準備を\n整えます", ACCENT),
]

for i, (step_no, title, desc, color) in enumerate(steps):
    x = Inches(0.5 + i * 3.15)
    y = Inches(1.6)

    # ステップ番号バー
    bar = add_rounded_rect(slide, x, y, Inches(2.9), Inches(0.6), color)
    set_text(bar, step_no, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # タイトル
    add_textbox(slide, x, y + Inches(0.75), Inches(2.9), Inches(0.5),
                title, font_size=18, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

    # 説明カード
    card = add_rounded_rect(slide, x, y + Inches(1.4), Inches(2.9), Inches(2.8), LIGHT_GRAY)
    add_textbox(slide, x + Inches(0.2), y + Inches(1.6), Inches(2.5), Inches(2.4),
                desc, font_size=13, color=DARK_GRAY)

    # 矢印（最後以外）
    if i < 3:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       x + Inches(3.0), y + Inches(2.5), Inches(0.3), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MID_GRAY
        arrow.line.fill.background()

# 下部メッセージ
add_rounded_rect(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.8), NAVY)
add_textbox(slide, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.6),
            "まずはSTEP 1のDB共有からスタート。3月中にエリア設計を完了し、4月のテスト配布に繋げます。",
            font_size=15, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════
# Slide 10: Thank You
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

add_rect(slide, Inches(0.8), Inches(2.8), Inches(1.5), Inches(0.06), BLUE)

add_multi_text(slide, Inches(0.8), Inches(3.0), Inches(11), Inches(2.5), [
    ("Thank you", 44, WHITE, True, PP_ALIGN.LEFT),
    ("", 12, WHITE, False, PP_ALIGN.LEFT),
    ("ご不明な点がございましたら、お気軽にお問い合わせください。", 18, RGBColor(0xAA, 0xCC, 0xEE), False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("株式会社アドレア", 16, RGBColor(0x99, 0xAA, 0xBB), False, PP_ALIGN.LEFT),
])

add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), BLUE)


# ─── 保存 ───
output_path = os.path.join(os.path.dirname(__file__), "proposal_v3.pptx")
prs.save(output_path)
print(f"保存完了: {output_path}")
print(f"スライド数: {len(prs.slides)}")
